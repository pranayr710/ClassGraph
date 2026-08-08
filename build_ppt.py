"""Build the ClassGraph Review-2 presentation.

Deck is structured 1:1 against the five evaluation criteria. Every number in
the results/problems slides is a real measured value from this repository's
own test runs -- see CHALLENGES_AND_SOLUTIONS.md for provenance.

Run:  python build_ppt.py
Out:  ClassGraph_Review2.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "ppt_assets"
OUT = ROOT / "ClassGraph_Review2.pptx"

# --------------------------------------------------------------------------- #
# Design system
# --------------------------------------------------------------------------- #

SW, SH = 13.333, 7.5          # slide size, inches (16:9)
ML = 0.68                      # left margin
MR = 0.68                      # right margin
CW = SW - ML - MR              # content width

INK = RGBColor(0x10, 0x27, 0x3F)        # deep navy - headings
BODY = RGBColor(0x3C, 0x50, 0x66)       # slate - body text
MUTE = RGBColor(0x7B, 0x8C, 0x9E)       # muted grey
TEAL = RGBColor(0x0E, 0x7C, 0x86)       # primary accent
TEAL_D = RGBColor(0x08, 0x59, 0x61)
AMBER = RGBColor(0xB4, 0x7A, 0x14)      # attention / problem
GREEN = RGBColor(0x1F, 0x6F, 0x50)      # solved / done
RED = RGBColor(0xA9, 0x33, 0x2A)        # blocker
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF3, 0xF7, 0xFA)      # card fill
PANEL2 = RGBColor(0xE9, 0xF1, 0xF4)     # tinted card
BORDER = RGBColor(0xDA, 0xE3, 0xEB)
RULE = RGBColor(0xE6, 0xEC, 0xF2)

FONT = "Segoe UI"
FONT_SB = "Segoe UI Semibold"
MONO = "Consolas"


def P(text, size=12, bold=False, color=BODY, font=FONT, italic=False,
      align=None, space_before=None, space_after=None, line=None):
    """One paragraph with a single run."""
    return {
        "runs": [{"t": text, "size": size, "bold": bold, "color": color,
                  "font": font, "italic": italic}],
        "align": align, "space_before": space_before,
        "space_after": space_after, "line": line,
    }


def PR(runs, align=None, space_before=None, space_after=None, line=None):
    """One paragraph with multiple runs. `runs` = list of dicts."""
    return {"runs": runs, "align": align, "space_before": space_before,
            "space_after": space_after, "line": line}


def R(t, size=12, bold=False, color=BODY, font=FONT, italic=False):
    return {"t": t, "size": size, "bold": bold, "color": color,
            "font": font, "italic": italic}


def add_text(slide, x, y, w, h, blocks, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, blk in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if blk.get("align") is not None:
            p.alignment = blk["align"]
        if blk.get("space_before") is not None:
            p.space_before = Pt(blk["space_before"])
        if blk.get("space_after") is not None:
            p.space_after = Pt(blk["space_after"])
        if blk.get("line") is not None:
            p.line_spacing = blk["line"]
        for seg in blk["runs"]:
            r = p.add_run()
            r.text = seg["t"]
            f = r.font
            f.name = seg.get("font", FONT)
            f.size = Pt(seg.get("size", 12))
            f.bold = seg.get("bold", False)
            f.italic = seg.get("italic", False)
            f.color.rgb = seg.get("color", BODY)
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=BORDER, line_w=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.045):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = adj
        except (IndexError, KeyError):
            pass
    s.text_frame.word_wrap = True
    return s


def bar(slide, x, y, w, h, fill):
    """Flat rectangle, no line -- for rules and accent bars."""
    return rect(slide, x, y, w, h, fill=fill, line=None, shape=MSO_SHAPE.RECTANGLE)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _wrapped_height(text, size, width_in, line=1.26):
    """Estimate rendered height so headers never silently sit on the content
    below them. Segoe UI averages ~0.50em glyph advance for mixed-case text."""
    chars_per_line = max(1.0, (width_in * 72.0) / (0.50 * size))
    import math as _m
    lines = max(1, _m.ceil(len(text) / chars_per_line))
    return lines * size * line / 72.0


def chrome(prs, criterion, title, lead=None, page=None):
    """Standard slide header + footer. Returns (slide, content_top).

    Title and lead boxes are sized from their own estimated wrapped height, so
    a two-line title pushes the content down instead of overlapping it.
    """
    s = new_slide(prs)
    y = 0.46
    if criterion:
        bar(s, ML, y + 0.045, 0.11, 0.16, TEAL)
        add_text(s, ML + 0.24, y, 8.0, 0.22,
                 [P(criterion.upper(), 10.5, True, TEAL, FONT_SB)])
        y += 0.30
    th = _wrapped_height(title, 27, CW) + 0.06
    add_text(s, ML, y, CW, th, [P(title, 27, True, INK, FONT_SB)])
    y += th + 0.11
    if lead:
        lh = _wrapped_height(lead, 12.5, CW - 0.2, line=1.18) + 0.04
        add_text(s, ML, y, CW - 0.2, lh, [P(lead, 12.5, False, MUTE, line=1.15)])
        y += lh + 0.10
    if page is not None:
        footer(s, page)
    return s, y + 0.10


def footer(s, page):
    bar(s, ML, SH - 0.62, CW, 0.012, RULE)
    add_text(s, ML, SH - 0.50, 6.0, 0.24,
             [P("ClassGraph  ·  Review 2", 9, False, MUTE)])
    add_text(s, SW - MR - 1.2, SH - 0.50, 1.2, 0.24,
             [P(str(page), 9, True, MUTE, align=PP_ALIGN.RIGHT)])


def card(slide, x, y, w, h, heading=None, lines=None, accent=TEAL,
         fill=PANEL, heading_size=13, body_size=10.5, accent_bar=True):
    """Panel with an optional left accent bar, heading and body lines."""
    rect(slide, x, y, w, h, fill=fill)
    if accent_bar:
        bar(slide, x, y + 0.10, 0.055, h - 0.20, accent)
    tx = x + (0.26 if accent_bar else 0.20)
    tw = w - (tx - x) - 0.20
    blocks = []
    if heading:
        blocks.append(P(heading, heading_size, True, INK, FONT_SB,
                        space_after=5, line=1.05))
    for ln in (lines or []):
        if isinstance(ln, str):
            blocks.append(P(ln, body_size, False, BODY, space_after=3, line=1.18))
        else:
            blocks.append(ln)
    if blocks:
        add_text(slide, tx, y + 0.16, tw, h - 0.30, blocks)


def stat(slide, x, y, w, h, value, label, color=TEAL):
    rect(slide, x, y, w, h, fill=PANEL)
    add_text(slide, x + 0.14, y + 0.16, w - 0.28, 0.60,
             [P(value, 30, True, color, FONT_SB, align=PP_ALIGN.CENTER)])
    add_text(slide, x + 0.14, y + 0.82, w - 0.28, h - 0.95,
             [P(label, 9.5, False, BODY, align=PP_ALIGN.CENTER, line=1.12)])


def table(slide, x, y, w, col_w, data, row_h=0.34, head_h=0.36,
          head_fill=INK, head_color=WHITE, size=9.5, head_size=9.5,
          zebra=True, col_colors=None, col_bold=None):
    """Simple, fully-styled table. data[0] is the header row."""
    rows, cols = len(data), len(data[0])
    total_h = head_h + row_h * (rows - 1)
    gt = slide.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                Inches(w), Inches(total_h))
    tbl = gt.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, cwid in enumerate(col_w):
        tbl.columns[i].width = Inches(cwid)
    for r in range(rows):
        tbl.rows[r].height = Inches(head_h if r == 0 else row_h)
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = head_fill
            else:
                cell.fill.fore_color.rgb = (
                    WHITE if (not zebra or r % 2 == 1) else PANEL
                )
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            f = run.font
            f.name = FONT_SB if r == 0 else FONT
            f.size = Pt(head_size if r == 0 else size)
            f.bold = bool(r == 0 or (col_bold and c in col_bold))
            if r == 0:
                f.color.rgb = head_color
            elif col_colors and c in col_colors:
                f.color.rgb = col_colors[c]
            else:
                f.color.rgb = BODY
    return tbl


def pill(slide, x, y, w, h, text, fill, color=WHITE, size=9):
    rect(slide, x, y, w, h, fill=fill, line=None, adj=0.5)
    add_text(slide, x, y + 0.035, w, h - 0.05,
             [P(text, size, True, color, FONT_SB, align=PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #

def s01_title(prs):
    s = new_slide(prs)
    bar(s, 0, 0, SW, 0.10, TEAL)
    bar(s, 0, SH - 0.055, SW, 0.055, INK)

    add_text(s, ML, 1.62, CW, 0.4,
             [P("CV PROJECT REVIEW 2", 11.5, True, TEAL, FONT_SB)])
    add_text(s, ML, 2.05, CW, 1.05, [P("ClassGraph", 58, True, INK, FONT_SB)])
    add_text(s, ML, 3.16, 9.6, 0.62,
             [P("Temporal Scene-Graph Group-Activity Analytics for Classrooms",
                20, False, TEAL_D, FONT_SB, line=1.1)])
    bar(s, ML, 3.95, 1.5, 0.028, TEAL)
    add_text(s, ML, 4.20, 9.3, 1.0,
             [P("Identity-persistent, behaviour-level engagement analytics — turning raw "
                "classroom video into per-student attention trajectories and group-level "
                "insight, without reducing students to a single score.",
                13, False, BODY, line=1.32)])

    labels = [("Perception + Identity", "Stages 1–2 built & tested"),
              ("Scene Graph + Temporal", "Stages 3–4 designed"),
              ("Group Activity", "Stage 5 designed")]
    bw = (CW - 0.44) / 3
    for i, (h, sub) in enumerate(labels):
        x = ML + i * (bw + 0.22)
        card(s, x, 5.55, bw, 0.86, heading=h, lines=[sub],
             heading_size=11.5, body_size=9.5, accent=TEAL if i == 0 else MUTE)
    return s


def s02_roadmap(prs, page):
    s, y = chrome(prs, None, "Review Roadmap",
                  "This deck is organised 1:1 against the five evaluation criteria.",
                  page)
    items = [
        ("1", "Problem Understanding", "Clarity of problem, objectives, motivation", 3),
        ("2", "Research Paper Selection", "Relevance, novelty, correctness of base paper", 6),
        ("3", "Literature Review Depth", "Related work identification & comparison", 7),
        ("4", "Dataset Identification", "Appropriateness, data format, challenges", 10),
        ("5", "Project Feasibility & Plan", "Workflow, timeline, results, risks", 12),
    ]
    lw = 7.35
    for i, (num, title, desc, slide_no) in enumerate(items):
        yy = y + 0.06 + i * 0.80
        rect(s, ML, yy, lw, 0.68, fill=WHITE, line=BORDER)
        rect(s, ML + 0.14, yy + 0.13, 0.42, 0.42, fill=INK, line=None, adj=0.22)
        add_text(s, ML + 0.14, yy + 0.19, 0.42, 0.30,
                 [P(num, 13, True, WHITE, FONT_SB, align=PP_ALIGN.CENTER)])
        add_text(s, ML + 0.70, yy + 0.11, 4.0, 0.26,
                 [P(title, 12.5, True, INK, FONT_SB)])
        add_text(s, ML + 0.70, yy + 0.36, 5.2, 0.24,
                 [P(desc, 9.5, False, MUTE)])
        add_text(s, ML + lw - 1.15, yy + 0.22, 1.0, 0.26,
                 [P(f"slide {slide_no}", 9, True, TEAL, align=PP_ALIGN.RIGHT)])

    px = ML + lw + 0.30
    pw = CW - lw - 0.30
    rect(s, px, y + 0.06, pw, 3.34, fill=INK, line=None)
    add_text(s, px + 0.30, y + 0.32, pw - 0.60, 0.3,
             [P("THE THESIS", 10, True, RGBColor(0x6F, 0xC5, 0xCE), FONT_SB)])
    add_text(s, px + 0.30, y + 0.72, pw - 0.60, 2.4,
             [P("Represent a classroom as a dynamic scene graph — students are nodes, "
                "attention behaviours are edges — keep identity stable within the session, "
                "and reason over how that graph evolves to recognise individual and group "
                "engagement.", 13.5, False, WHITE, line=1.34)])
    card(s, px, y + 3.56, pw, 0.86,
         heading="Status at this review",
         lines=["Stages 1–2 implemented, 109 automated tests passing, "
                "validated on real classroom footage."],
         heading_size=11.5, body_size=9.5, accent=GREEN, fill=PANEL2)
    return s


def s03_problem(prs, page):
    s, y = chrome(prs, "Criterion 1 — Problem Understanding",
                  "Attendance is solved. Engagement is not.",
                  "The classroom already has cameras. What it does not have is any way to "
                  "answer the question a teacher actually cares about.", page)

    card(s, ML, y, 6.15, 1.62, heading="The problem",
         lines=["A teacher cannot continuously observe 30–40 students while also teaching.",
                "Manual observation is subjective, sampled at a few moments, and does not "
                "scale to every student across a whole lecture."],
         accent=RED)
    card(s, ML + 6.45, y, 6.15, 1.62, heading="Why existing tools fall short",
         lines=["Attendance and face-recognition systems answer “who is present”, "
                "not “who is engaged”.",
                "Emotion classifiers give a per-frame label with no memory — no system "
                "reasons about behaviour over time or across the group."],
         accent=RED)

    yy = y + 1.86
    add_text(s, ML, yy, CW, 0.3,
             [P("What we mean by “engagement” — defined before it is measured",
                13.5, True, INK, FONT_SB)])
    add_text(s, ML, yy + 0.34, CW, 1.00,
             [PR([R("We measure ", 11), R("observable behaviour", 11, True, TEAL),
                  R(", not inner emotional state. ", 11),
                  R("Where a student is looking, whether their eyes are open, whether a "
                    "phone is present, how their body is oriented, and how all of that "
                    "changes over time. ", 11),
                  R("We never infer feelings — that is both scientifically weaker and, in "
                    "education settings, legally restricted.", 11, False, AMBER)],
                 line=1.22)])

    yy += 1.46
    objectives = [
        ("Detect", "every student in frame, including back rows"),
        ("Track", "stable identity for the whole session"),
        ("Estimate", "attention state over a time window"),
        ("Flag", "device distraction & eye closure"),
        ("Relate", "peer orientation & group activity"),
        ("Report", "trajectories + class-level trends"),
    ]
    bw = (CW - 5 * 0.16) / 6
    for i, (verb, desc) in enumerate(objectives):
        x = ML + i * (bw + 0.16)
        rect(s, x, yy, bw, 1.12, fill=PANEL2, line=BORDER)
        add_text(s, x + 0.14, yy + 0.14, bw - 0.28, 0.26,
                 [P(verb, 12, True, TEAL_D, FONT_SB)])
        add_text(s, x + 0.14, yy + 0.44, bw - 0.28, 0.62,
                 [P(desc, 9, False, BODY, line=1.16)])
    return s


def s04_motivation(prs, page):
    s, y = chrome(prs, "Criterion 1 — Motivation",
                  "Why this matters, and why now",
                  None, page)
    stats = [("30–40", "students one teacher must monitor simultaneously in a typical class"),
             ("~50%", "of detected students have no usable face from a real ceiling camera — "
                      "measured on our own footage"),
             ("5", "mature CV subfields that can now be combined without training from scratch"),
             ("0", "deployed tools that track engagement per student over a full session")]
    bw = (CW - 3 * 0.24) / 4
    for i, (v, lab) in enumerate(stats):
        stat(s, ML + i * (bw + 0.24), y, bw, 1.62, v, lab,
             color=TEAL if i in (0, 2) else AMBER)

    yy = y + 1.90
    card(s, ML, yy, 6.15, 1.42, heading="Why now — feasibility changed",
         lines=["Detection, face mesh, head pose and body pose are all available as strong "
                "pretrained models. The research problem is no longer “can we see it” but "
                "“what does it mean, and over what timescale”.",
                ], accent=TEAL)
    card(s, ML + 6.45, yy, 6.15, 1.42, heading="Where it applies",
         lines=["Schools, colleges and coaching centres · hybrid learning · corporate "
                "training · teaching-quality feedback · educational research on attention.",
                ], accent=TEAL)

    yy += 1.66
    rect(s, ML, yy, CW, 1.14, fill=INK, line=None)
    add_text(s, ML + 0.30, yy + 0.16, CW - 0.60, 0.88,
             [PR([R("Our framing:  ", 12, True, RGBColor(0x6F, 0xC5, 0xCE), FONT_SB),
                  R("the useful output is not a score per student — it is a ", 12, False, WHITE),
                  R("trend a teacher can act on", 12, True, WHITE, FONT_SB),
                  R(". “Attention across the room dropped sharply 22 minutes in” is "
                    "actionable and fair. “Student 7 is disengaged” is neither.",
                    12, False, WHITE)], line=1.26)])
    return s


def s05_objectives(prs, page):
    s, y = chrome(prs, "Criterion 1 — Objectives",
                  "Objectives, each with a measurable success criterion",
                  "Stated so progress can be checked rather than asserted. "
                  "Status reflects this review.", page)
    data = [
        ["#", "Objective", "Measurable success criterion", "Status"],
        ["O1", "Detect every student in frame",
         "Recall on real classroom frames; back rows not lost to downscaling", "Done"],
        ["O2", "Maintain stable identity within a session",
         "One track_id per student across occlusion; identity never crosses sessions", "Done"],
        ["O3", "Per-student attention state over time",
         "Windowed distribution over 6 behaviour categories, not per-frame verdicts", "Done"],
        ["O4", "Flag device distraction and eye closure",
         "Phone-overlap + gaze-down evidence; EAR-based eye state", "Done"],
        ["O5", "Model peer orientation and group activity",
         "Geometric pair detection now; scene-graph + GAR next", "In progress"],
        ["O6", "Report trajectories and class-level trends",
         "Class summary by default, individual as explicit drill-down", "Partial"],
    ]
    tbl = table(s, ML, y, CW, [0.52, 3.30, 6.55, 1.60], data,
                row_h=0.52, head_h=0.40, size=10, head_size=10,
                col_bold={1})
    # colour the status column by value
    status_colors = {"Done": GREEN, "In progress": AMBER, "Partial": AMBER}
    for r in range(1, len(data)):
        cell = tbl.cell(r, 3)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = status_colors.get(data[r][3], BODY)
        run.font.bold = True

    add_text(s, ML, y + 3.60, CW, 0.5,
             [PR([R("Note on scope:  ", 10.5, True, INK, FONT_SB),
                  R("O1–O4 are implemented and covered by automated tests. O5–O6 are "
                    "partially built — the honest position at this review is a working "
                    "perception + attention core with the relational layer designed and "
                    "started, not finished.", 10.5)], line=1.2)])
    return s


def s06_papers(prs, page):
    s, y = chrome(prs, "Criterion 2 — Research Paper Selection",
                  "Three base papers, one per technical pillar",
                  "All three are peer-reviewed at top-tier venues (CVPR / ICCV), have public "
                  "code, and map 1:1 onto a pillar of our architecture.", page)
    papers = [
        (TEAL, "PILLAR 1 · IDENTITY",
         "Clothes-Changing Person Re-Identification With RGB Modality Only",
         "Gu, Chang, Ma, Bai, Shan & Chen — CVPR 2022, pp. 1060–1069",
         "Proposes a Clothes-based Adversarial Loss (CAL) that penalises the model's "
         "ability to predict clothing, forcing it to learn clothes-irrelevant identity "
         "cues (face, hairstyle, body shape, gait) from RGB alone.",
         "What we take: identity that survives appearance change and occlusion, so a "
         "student who turns away and back is not counted as a new person."),
        (TEAL, "PILLAR 2 · RELATIONS OVER TIME",
         "Spatial-Temporal Transformer for Dynamic Scene Graph Generation",
         "Cong, Liao, Ackermann, Rosenhahn & Yang — ICCV 2021",
         "STTran pairs a spatial encoder (relations within a frame) with a temporal "
         "decoder (dependencies across frames), and accepts variable-length video "
         "without clipping. Benchmarked on Action Genome.",
         "What we take: the dynamic scene-graph formulation and the temporal decoder — "
         "our classroom graph must evolve, not be re-guessed every frame."),
        (TEAL, "PILLAR 3 · THE GROUP",
         "Learning Actor Relation Graphs for Group Activity Recognition",
         "Wu, Wang, Wang, Guo & Wu — CVPR 2019",
         "Builds an Actor Relation Graph capturing both appearance and position "
         "relations between actors, learned end-to-end with a GCN; adds spatially "
         "localised and temporally randomised sparsification. Benchmarked on "
         "Volleyball and Collective Activity.",
         "What we take: the actor-relation graph + GCN readout to turn per-student "
         "states into one class-level activity label."),
    ]
    cwid = (CW - 2 * 0.24) / 3
    for i, (accent, tag, title, cite, method, take) in enumerate(papers):
        x = ML + i * (cwid + 0.24)
        rect(s, x, y, cwid, 3.92, fill=WHITE, line=BORDER)
        bar(s, x, y, cwid, 0.055, accent)
        add_text(s, x + 0.22, y + 0.24, cwid - 0.44, 0.24,
                 [P(tag, 9, True, TEAL, FONT_SB)])
        add_text(s, x + 0.22, y + 0.54, cwid - 0.44, 0.90,
                 [P(title, 12.5, True, INK, FONT_SB, line=1.10)])
        add_text(s, x + 0.22, y + 1.48, cwid - 0.44, 0.42,
                 [P(cite, 9, True, TEAL_D, line=1.16)])
        add_text(s, x + 0.22, y + 1.96, cwid - 0.44, 1.30,
                 [P(method, 9.5, False, BODY, line=1.20)])
        bar(s, x + 0.22, y + 3.26, cwid - 0.44, 0.012, RULE)
        add_text(s, x + 0.22, y + 3.38, cwid - 0.44, 0.48,
                 [P(take, 9, True, TEAL_D, line=1.18)])

    rect(s, ML, y + 4.10, CW, 0.80, fill=PANEL2, line=BORDER)
    add_text(s, ML + 0.26, y + 4.24, CW - 0.52, 0.56,
             [PR([R("Correctness of selection:  ", 11, True, INK, FONT_SB),
                  R("each paper is the standard reference for its pillar, not a peripheral "
                    "choice. Our novelty is not in re-deriving any one of them — it is in "
                    "fusing all three in a domain none of them targets.", 11)], line=1.22)])
    return s


def s07_lit_technical(prs, page):
    s, y = chrome(prs, "Criterion 3 — Literature Review Depth",
                  "Related work, method, result, and the gap each leaves open",
                  None, page)
    data = [
        ["Area", "Key work (venue, year)", "Method in one line", "Benchmark", "Limitation for our setting"],
        ["Cloth-changing\nRe-ID",
         "CAL — Gu et al.\n(CVPR 2022)",
         "Adversarial loss suppresses clothing cues; RGB only",
         "PRCC, LTCC,\nDeepChange",
         "Built for surveillance re-identification; never applied to session-scoped classroom identity"],
        ["", "AIM (CVPR 2023);\nDeepChange (2023)",
         "Causal/attribute-based identity disentangling; large cloth-change benchmark",
         "PRCC, LTCC",
         "Same domain mismatch; no notion of a fixed seating layout to exploit"],
        ["Video scene\ngraph",
         "STTran — Cong et al.\n(ICCV 2021)",
         "Spatial encoder + temporal decoder over variable-length video",
         "Action Genome",
         "Generic object–object predicates; nothing for attention, gaze or education"],
        ["Group activity\nrecognition",
         "ARG — Wu et al.\n(CVPR 2019)",
         "Actor relation graph (appearance + position) with GCN",
         "Volleyball,\nCollective Activity",
         "Sports/crowd domains; group label per clip, not sustained engagement"],
        ["", "GroupFormer\n(ICCV 2021)",
         "Joint spatial-temporal transformer with clustered attention",
         "Volleyball,\nCollective Activity",
         "Stronger backbone, same domain gap; assumes clean actor tracks"],
        ["Classroom CV",
         "SCB-Dataset\n(2023–2025)",
         "Large labelled classroom behaviour detection benchmark (YOLO-based)",
         "SCB (own)",
         "Per-frame detection only — no identity across time, no relational or group reasoning"],
        ["", "DAiSEE; EmotiW\nengagement",
         "Affective engagement classification from webcam video",
         "DAiSEE",
         "Single-subject webcam framing and affect labels — not a multi-student room"],
    ]
    table(s, ML, y, CW, [1.28, 2.05, 3.40, 1.55, 3.69], data,
          row_h=0.62, head_h=0.36, size=8.5, head_size=9, col_bold={0})

    add_text(s, ML, y + 4.72, CW, 0.44,
             [PR([R("Reading across the rows:  ", 10.5, True, INK, FONT_SB),
                  R("identity, relations and group reasoning have each been solved "
                    "separately, and classroom CV has solved per-frame detection. "
                    "No row does two of these at once.", 10.5)], line=1.2)])
    return s


def s08_lit_behavioural(prs, page):
    s, y = chrome(prs, "Criterion 3 — Literature Review Depth",
                  "The second half: what should the labels even be?",
                  "A CV-only review would stop at the previous slide. But the hardest question "
                  "here is not architectural — it is what counts as “engaged”, and over what "
                  "timescale. We reviewed the education and psychology literature to answer it.",
                  page)
    left = [
        ("The problem with the field",
         ["Khan, Abedi & Colella (2022/23) reviewed engagement-detection systems and found "
          "most invent their own ad-hoc label scheme instead of grounding it in a validated "
          "instrument — making results incomparable.",
          "We treated that as a warning about our own design, not a citation to drop in."],
         RED),
        ("So we grounded the taxonomy",
         ["BOSS (Shapiro; reviewed in Volpe et al., 2005) is a validated, individual-student "
          "classroom observation instrument: Active/Passive Engaged Time vs Off-task "
          "Motor/Verbal/Passive.",
          "Our six categories map onto it — and where they cannot, we say so."],
         GREEN),
    ]
    right = [
        ("And grounded the timescale",
         ["Ariga & Lleras (2011): a single brief break restored vigilance over a 50-minute "
          "task — micro-lapses are the mechanism that protects attention, not evidence of "
          "its failure.",
          "Faber, Bixler & D'Mello: gaze-based mind-wandering detection works best over "
          "~12-second windows.",
          "Zakszeski et al. (2017): 5–15 s momentary sampling tracks continuous observation "
          "most accurately."],
         TEAL),
        ("And bounded what vision can claim",
         ["Kendon's F-formations give a geometric definition of joint orientation — usable "
          "from body pose alone.",
          "But Bassiou et al. (Interspeech 2016) show that separating productive academic "
          "peer talk from idle chat needs audio. We therefore refuse to make that call."],
         AMBER),
    ]
    colw = (CW - 0.30) / 2
    for col, (items, x) in enumerate([(left, ML), (right, ML + colw + 0.30)]):
        yy = y
        for heading, lines, accent in items:
            h = 1.42 if col == 0 else 1.72
            card(s, x, yy, colw, h, heading=heading, lines=lines,
                 accent=accent, heading_size=12, body_size=9.5)
            yy += h + 0.20

    rect(s, ML, y + 3.68, CW, 0.86, fill=INK, line=None)
    add_text(s, ML + 0.28, y + 3.82, CW - 0.56, 0.60,
             [PR([R("Why this is depth, not decoration:  ", 11, True,
                    RGBColor(0x6F, 0xC5, 0xCE), FONT_SB),
                  R("the 15-second window in our implementation was chosen from this "
                    "literature, and it is independently corroborated by two unrelated "
                    "lines of work (gaze-based mind-wandering, and school-psychology "
                    "sampling methodology).", 11, False, WHITE)], line=1.24)])
    return s


def s09_gap(prs, page):
    s, y = chrome(prs, "Criterion 3 — Research Gap",
                  "The white space, stated precisely",
                  None, page)
    cols = [
        ("Identity persists", "Re-ID gives a stable student across a session",
         "but only ever demonstrated for surveillance"),
        ("Relations are modelled", "VidSGG gives evolving subject–predicate–object graphs",
         "but with generic predicates, never attention"),
        ("The group is read out", "GAR gives one activity label for a set of actors",
         "but on sports clips, not sustained engagement"),
    ]
    cwid = (CW - 2 * 0.26) / 3
    for i, (h, has, lacks) in enumerate(cols):
        x = ML + i * (cwid + 0.26)
        rect(s, x, y, cwid, 1.86, fill=WHITE, line=BORDER)
        add_text(s, x + 0.22, y + 0.22, cwid - 0.44, 0.3,
                 [P(h, 13, True, INK, FONT_SB)])
        add_text(s, x + 0.22, y + 0.62, cwid - 0.44, 0.56,
                 [PR([R("✓ ", 11, True, GREEN), R(has, 10, False, BODY)], line=1.18)])
        add_text(s, x + 0.22, y + 1.24, cwid - 0.44, 0.52,
                 [PR([R("✗ ", 11, True, RED), R(lacks, 10, False, MUTE)], line=1.18)])

    yy = y + 2.14
    rect(s, ML, yy, CW, 1.30, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, ML + 0.30, yy + 0.20, CW - 0.60, 0.30,
             [P("THE GAP WE ADDRESS", 10, True, TEAL_D, FONT_SB)])
    add_text(s, ML + 0.30, yy + 0.56, CW - 0.60, 0.62,
             [P("No prior work combines identity persistence, relational scene graphs and "
                "temporal group reasoning for classrooms — and none of the classroom systems "
                "ground their behaviour labels in a validated observation instrument.",
                14, False, INK, line=1.26)])

    yy += 1.58
    card(s, ML, yy, CW, 1.06,
         heading="Our claim is deliberately narrow",
         lines=["We are not claiming a new Re-ID loss, a new scene-graph architecture, or a "
                "new GAR backbone. The contribution is the integration and the domain: an "
                "identity-persistent, behaviour-grounded, temporally-reasoned classroom "
                "analytics system — plus the honest reporting of where vision alone cannot "
                "decide, which the reviewed classroom literature consistently glosses over."],
         accent=TEAL, heading_size=12, body_size=10.5)
    return s


def s10_dataset_primary(prs, page):
    s, y = chrome(prs, "Criterion 4 — Dataset Identification",
                  "Primary dataset: SCB-Dataset (Student Classroom Behavior)",
                  "Selected after comparing classroom-specific options against our actual "
                  "deployment conditions and our actual failure modes.", page)

    card(s, ML, y, 4.05, 3.30, heading="Why this one",
         lines=["Real classroom images from elevated / rear-corner cameras — the same "
                "viewpoint that causes our hardest problems, not clean frontal webcam video.",
                "Covers both student and teacher behaviour in one label space.",
                "Benchmarked with YOLO-series detectors, so it drops directly into our "
                "existing YOLOv11 detection stage.",
                "Public and actively maintained (arXiv 2304.02488, v7 · Aug 2025)."],
         accent=GREEN, heading_size=13, body_size=10)

    x2 = ML + 4.05 + 0.26
    rect(s, x2, y, 3.55, 3.30, fill=PANEL, line=BORDER)
    add_text(s, x2 + 0.24, y + 0.20, 3.1, 0.3,
             [P("Scale & format", 13, True, INK, FONT_SB)])
    rows = [("Detection split", "13,330 images"),
            ("Detection labels", "122,977 boxes"),
            ("Classification split", "21,019 images"),
            ("Behaviour classes", "19–20"),
            ("Annotation format", "bbox + class / image"),
            ("Benchmarked with", "YOLO series, VLMs")]
    for i, (k, v) in enumerate(rows):
        yy = y + 0.62 + i * 0.42
        add_text(s, x2 + 0.24, yy, 1.85, 0.28, [P(k, 9.5, False, MUTE)])
        add_text(s, x2 + 2.05, yy, 1.28, 0.28,
                 [P(v, 9.5, True, INK, MONO, align=PP_ALIGN.RIGHT)])

    x3 = x2 + 3.55 + 0.26
    w3 = SW - MR - x3
    rect(s, x3, y, w3, 3.30, fill=INK, line=None)
    add_text(s, x3 + 0.26, y + 0.20, w3 - 0.52, 0.3,
             [P("THE DECIDING FACTOR", 10, True, RGBColor(0x6F, 0xC5, 0xCE), FONT_SB)])
    add_text(s, x3 + 0.26, y + 0.58, w3 - 0.52, 0.80,
             [P("Its label set resolves the exact ambiguity our geometry cannot.",
                13.5, True, WHITE, FONT_SB, line=1.16)])
    add_text(s, x3 + 0.26, y + 1.44, w3 - 0.52, 1.70,
             [P("From head pose alone we can tell a student is looking down — but not "
                "whether they are reading, writing, or disengaged. SCB labels "
                "bow head, read, write, using the phone, turn head, talk and discuss "
                "as separate classes, which is supervision for precisely the distinction "
                "we documented as unresolvable without it.",
                10.5, False, RGBColor(0xD6, 0xE4, 0xEC), line=1.26)])

    yy = y + 3.56
    add_text(s, ML, yy, CW, 0.28,
             [P("Class labels that map directly onto our behaviour categories",
                11.5, True, INK, FONT_SB)])
    tags = ["hand-raising", "read", "write", "bow head", "turn head", "talk",
            "discuss", "using the phone", "using the computer", "stand",
            "answer", "yawn", "leaning on desk", "teacher"]
    x = ML
    yy += 0.36
    for t in tags:
        w = 0.28 + 0.088 * len(t)
        if x + w > SW - MR:
            x = ML
            yy += 0.40
        pill(s, x, yy, w, 0.32, t, PANEL2, TEAL_D, size=9)
        x += w + 0.12
    return s


def s11_dataset_support(prs, page):
    s, y = chrome(prs, "Criterion 4 — Dataset Identification",
                  "Supporting benchmarks per module, and the honest challenges",
                  "No dataset carries classroom video + identity + scene-graph + group labels "
                  "together, so each module is benchmarked on its field-standard set.", page)
    data = [
        ["Module", "Dataset", "Label / format", "Why this one"],
        ["Person + object detection", "SCB-Dataset · COCO",
         "bbox + class per image", "Classroom-native; COCO covers phone/laptop/book"],
        ["Identity (Re-ID)", "PRCC · LTCC · DeepChange",
         "image + identity label, clothing change", "Standard cloth-changing Re-ID protocol (CAL's own benchmarks)"],
        ["Scene graph", "Action Genome",
         "subject–predicate–object per frame", "STTran's benchmark — lets us compare against the base paper"],
        ["Group activity", "Volleyball · Collective Activity",
         "clip + group activity label", "ARG's benchmarks — same reason"],
        ["Head pose", "300W-LP · BIWI · AFLW2000",
         "image + yaw/pitch/roll", "Training/eval sets of the SixDRepNet model we deploy"],
        ["Engagement (reference)", "DAiSEE",
         "video + engagement level", "Comparison point for engagement, though single-subject"],
    ]
    table(s, ML, y, CW, [2.35, 2.55, 2.85, 4.22], data,
          row_h=0.42, head_h=0.36, size=9, head_size=9.5, col_bold={0})

    yy = y + 3.02
    challenges = [
        ("Fragmentation", "No dataset combines classroom + identity + group labels.",
         "Benchmark per module against its standard set; integrate on our own footage."),
        ("SCB has no identity labels", "It is per-frame detection only — Re-ID cannot be "
         "evaluated on it.", "Re-ID benchmarked separately on PRCC/LTCC; identity validated "
         "on our own continuous video."),
        ("Resolution sensitivity", "Our own 13-image pilot showed face detection collapsing "
         "at low resolution.", "Measured it explicitly (see results) and set a minimum "
         "capture resolution as a requirement."),
    ]
    cwid = (CW - 2 * 0.24) / 3
    for i, (h, prob, mit) in enumerate(challenges):
        x = ML + i * (cwid + 0.24)
        rect(s, x, yy, cwid, 1.50, fill=WHITE, line=BORDER)
        bar(s, x, yy, cwid, 0.05, AMBER)
        add_text(s, x + 0.20, yy + 0.16, cwid - 0.40, 0.26,
                 [P(h, 11, True, INK, FONT_SB)])
        add_text(s, x + 0.20, yy + 0.46, cwid - 0.40, 0.52,
                 [P(prob, 9, False, BODY, line=1.16)])
        add_text(s, x + 0.20, yy + 1.02, cwid - 0.40, 0.44,
                 [PR([R("→ ", 9, True, GREEN), R(mit, 9, False, TEAL_D)], line=1.16)])

    add_text(s, ML, yy + 1.62, CW, 0.30,
             [PR([R("Planned contribution:  ", 10.5, True, INK, FONT_SB),
                  R("annotate 15–30 minutes of our own classroom video with identity + "
                    "attention + interaction labels — turning the fragmentation gap into a "
                    "small dataset contribution.", 10.5)], line=1.2)])
    return s


def s12_architecture(prs, page):
    s, y = chrome(prs, "Criterion 5 — Workflow",
                  "System architecture — five stages, frozen contracts",
                  "Each stage consumes the previous stage's output through a fixed JSON "
                  "schema, so three people can build in parallel against stable interfaces.",
                  page)
    stages = [
        ("1", "Perception", "YOLOv11 · MediaPipe Face Mesh\nMediaPipe Pose · SixDRepNet",
         "who & where", GREEN, "BUILT"),
        ("2", "Identity", "ByteTrack now\nCAL Re-ID next",
         "stable IDs", GREEN, "BUILT"),
        ("3", "Scene Graph", "nodes = students\nedges = behaviours",
         "what relates to what", MUTE, "NEXT"),
        ("4", "Temporal", "STTran-style decoder\nover the graph",
         "how it evolves", MUTE, "NEXT"),
        ("5", "Group Activity", "ARG-style graph\nreadout",
         "the class as a whole", MUTE, "NEXT"),
    ]
    gap = 0.30
    bw = (CW - 4 * gap) / 5
    for i, (num, title, tech, out, accent, badge) in enumerate(stages):
        x = ML + i * (bw + gap)
        done = badge == "BUILT"
        rect(s, x, y, bw, 2.72, fill=WHITE if done else PANEL, line=accent if done else BORDER,
             line_w=1.5 if done else 0.75)
        bar(s, x, y, bw, 0.05, accent)
        add_text(s, x + 0.18, y + 0.22, 0.4, 0.3, [P(num, 15, True, accent, FONT_SB)])
        pw = 0.68
        pill(s, x + bw - pw - 0.18, y + 0.24, pw, 0.26, badge,
             GREEN if done else RGBColor(0xC3, 0xCD, 0xD6), WHITE, size=7.5)
        add_text(s, x + 0.18, y + 0.62, bw - 0.36, 0.34,
                 [P(title, 12.5, True, INK, FONT_SB)])
        add_text(s, x + 0.18, y + 1.02, bw - 0.36, 0.86,
                 [P(tech, 9, False, BODY, MONO, line=1.24)])
        bar(s, x + 0.18, y + 2.00, bw - 0.36, 0.012, RULE)
        add_text(s, x + 0.18, y + 2.14, bw - 0.36, 0.44,
                 [P(out, 9.5, True, TEAL_D, line=1.14)])
        if i < 4:
            add_text(s, x + bw + 0.03, y + 1.14, gap - 0.06, 0.3,
                     [P("▸", 15, True, RGBColor(0xB6, 0xC4, 0xCF),
                        align=PP_ALIGN.CENTER)])

    yy = y + 3.00
    card(s, ML, yy, (CW - 0.26) / 2, 1.16,
         heading="The frozen contract",
         lines=["One JSON object per frame: persons (bbox, track_id, face landmarks + EAR, "
                "head pose + gaze label, body pose) and objects (phone/laptop/book). "
                "Validated with JSON Schema on every run — 321/321 records valid on our "
                "last full-video test."],
         accent=TEAL, heading_size=12, body_size=9.5)
    card(s, ML + (CW - 0.26) / 2 + 0.26, yy, (CW - 0.26) / 2, 1.16,
         heading="Privacy by design",
         lines=["Identity is computed from motion/geometry and scoped to a single session — "
                "there is no face-recognition database. Two regression tests enforce that "
                "identity cannot leak between sessions, which is the legal line between "
                "attention analytics and biometric surveillance in schools."],
         accent=GREEN, heading_size=12, body_size=9.5)
    return s


def s13_status(prs, page):
    s, y = chrome(prs, "Criterion 5 — Feasibility",
                  "What already runs, on our own hardware",
                  "Feasibility is not asserted here — it is demonstrated. The perception and "
                  "identity stages are complete, tested and validated on real footage.", page)
    stats = [("109", "automated tests passing\n(0 skipped, 0 failing)", GREEN),
             ("321/321", "JSONL records valid against\nthe frozen schema", GREEN),
             ("7.8", "FPS end-to-end on 4K video,\nRTX 4050 laptop GPU", TEAL),
             ("100%", "of detected students carry\nat least one usable signal", GREEN)]
    bw = (CW - 3 * 0.24) / 4
    for i, (v, lab, c) in enumerate(stats):
        stat(s, ML + i * (bw + 0.24), y, bw, 1.58, v, lab, color=c)

    yy = y + 1.86
    add_text(s, ML, yy, CW, 0.3,
             [P("Stage-by-stage status", 13.5, True, INK, FONT_SB)])
    yy += 0.40
    data = [
        ["Stage", "Component", "State", "Evidence"],
        ["1", "YOLOv11 person + object detection", "Complete",
         "236 persons across 12 real classroom images"],
        ["1", "MediaPipe Face Mesh + EAR", "Complete",
         "468 landmarks; EAR range 0.10–0.49 on real faces"],
        ["1", "SixDRepNet head pose + gaze label", "Complete",
         "5-class gaze label; sign convention verified against library source"],
        ["1", "MediaPipe Pose fallback", "Complete",
         "recovers 56% of students who have no detectable face"],
        ["2", "ByteTrack identity", "Complete",
         "stable track_id across occlusion; session-scoped by test"],
        ["3–5", "Scene graph · temporal · group activity", "Designed",
         "contracts defined; base papers selected and read"],
    ]
    tbl = table(s, ML, yy, CW, [0.72, 4.60, 1.55, 5.10], data,
                row_h=0.42, head_h=0.36, size=9.5, head_size=9.5, col_bold={1})
    for r in range(1, len(data)):
        run = tbl.cell(r, 2).text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = GREEN if data[r][2] == "Complete" else AMBER
        run.font.bold = True
    return s


def s14_results(prs, page):
    s, y = chrome(prs, "Criterion 5 — Results",
                  "Measured improvements, before and after",
                  "Every figure below is from a real run on real classroom images in this "
                  "repository — not an estimate.", page)
    data = [
        ["Metric (same 12-image real classroom set)", "Before", "After", "Change"],
        ["Persons detected", "139", "236", "+70%"],
        ["Faces detected on real footage", "0", "95", "0 → usable"],
        ["Students with any usable signal (321-frame video)", "265 / 321", "321 / 321", "100% coverage"],
        ["Gaze label “down” reachable at all", "No (bug)", "Yes", "fixed"],
        ["Automated tests passing", "10 (9 skipped)", "109 (0 skipped)", "+99"],
        ["Throughput, 4K video on RTX 4050", "11.0 FPS", "7.8 FPS", "cost of posture"],
    ]
    tw = 8.20
    tbl = table(s, ML, y, tw, [4.40, 1.25, 1.25, 1.30], data,
                row_h=0.42, head_h=0.38, size=9.5, head_size=9.5, col_bold={0})
    for r in range(1, len(data)):
        for c in (1, 2, 3):
            run = tbl.cell(r, c).text_frame.paragraphs[0].runs[0]
            run.font.name = MONO
            run.font.size = Pt(9)
            if c == 3:
                # last row is a cost, not a win -- colour it honestly
                run.font.color.rgb = AMBER if r == len(data) - 1 else GREEN
                run.font.bold = True
            elif c == 1:
                run.font.color.rgb = MUTE

    yy = y + 0.38 + 6 * 0.42 + 0.26
    card(s, ML, yy, tw, 1.06,
         heading="Why the throughput row is on this slide",
         lines=["Adding the body-pose fallback took us from 11.0 to 7.8 FPS. We are reporting "
                "the cost next to the benefit rather than quoting only the coverage gain — "
                "the trade is worth it for offline analysis, and it is the number a reviewer "
                "should be able to challenge us on."],
         accent=AMBER, heading_size=11.5, body_size=9.5, fill=PANEL)

    img_w = 3.40
    img_x = SW - MR - img_w
    a = ASSETS / "annot_baseline_960_c40.jpg"
    b = ASSETS / "annot_candidate_1536_c30.jpg"
    if a.exists() and b.exists():
        ih = img_w * 1088 / 1920
        add_text(s, img_x, y, img_w, 0.24,
                 [P("BEFORE — 960 px, conf 0.40 → 20 found", 9, True, RED, FONT_SB)])
        s.shapes.add_picture(str(a), Inches(img_x), Inches(y + 0.28),
                             width=Inches(img_w))
        y2 = y + 0.28 + ih + 0.20
        add_text(s, img_x, y2, img_w, 0.24,
                 [P("AFTER — 1536 px, conf 0.30 → 35 found", 9, True, GREEN, FONT_SB)])
        s.shapes.add_picture(str(b), Inches(img_x), Inches(y2 + 0.28),
                             width=Inches(img_w))
        add_text(s, img_x, y2 + 0.28 + ih + 0.10, img_w, 0.44,
                 [P("Rendered at 1536 px to make the effect visible; the shipped default is "
                    "1280 px, which captures most of the gain at 1.5× the cost.",
                    8, False, MUTE, line=1.16)])
    return s


def s15_problems_perception(prs, page):
    s, y = chrome(prs, "Criterion 5 — Problems Faced (1 / 2)",
                  "Perception problems, root causes, and what we did",
                  "These were found by running the pipeline on real footage — none of them "
                  "were visible in synthetic tests.", page)

    probs = [
        ("PROBLEM 1", "Camera angle — a hard ceiling, not a bug",
         "Real classroom cameras sit high in a rear corner. A student bowed over a desk "
         "shows the camera the crown of their head. We measured a ~45% ceiling on face "
         "availability: 130 of 236 detected students had no face that ANY face model "
         "could find.",
         "Accepted it as a constraint instead of chasing it. Added a face-independent "
         "signal (body pose) so those students are still measurable — coverage went from "
         "265/321 to 321/321.", RED),
        ("PROBLEM 2", "Face detection returned zero faces on real footage",
         "Face Mesh was run once over the whole frame. MediaPipe downscales its input, so "
         "a face that is small relative to a 4K frame is destroyed before detection runs. "
         "Result: 0 faces on every real image, while synthetic tests passed.",
         "Run Face Mesh per detected person crop instead. 0/139 → 95/139 faces. We also "
         "measured that padding the crop makes it worse, and shipped zero padding against "
         "our own initial intuition.", RED),
    ]
    LB = 8.40                      # left block: the two detailed problems
    img_w = CW - LB - 0.30         # right column: the fallback, working
    cwid = (LB - 0.28) / 2
    for i, (tag, title, cause, fix, accent) in enumerate(probs):
        x = ML + i * (cwid + 0.28)
        rect(s, x, y, cwid, 3.30, fill=WHITE, line=BORDER)
        bar(s, x, y, cwid, 0.05, accent)
        add_text(s, x + 0.24, y + 0.20, cwid - 0.48, 0.24,
                 [P(tag, 9, True, accent, FONT_SB)])
        add_text(s, x + 0.24, y + 0.48, cwid - 0.48, 0.52,
                 [P(title, 13, True, INK, FONT_SB, line=1.08)])
        add_text(s, x + 0.24, y + 1.04, cwid - 0.48, 1.12,
                 [P(cause, 10, False, BODY, line=1.20)])
        bar(s, x + 0.24, y + 2.22, cwid - 0.48, 0.012, RULE)
        add_text(s, x + 0.24, y + 2.34, cwid - 0.48, 0.92,
                 [PR([R("Solution:  ", 10, True, GREEN, FONT_SB),
                      R(fix, 10, False, TEAL_D)], line=1.20)])

    # The posture fallback is the answer to Problem 1 -- show it working.
    pose = ASSETS / "t1_pose.jpg"
    if pose.exists():
        px = ML + LB + 0.30
        add_text(s, px, y, img_w, 0.24,
                 [P("THE FALLBACK, WORKING", 8.5, True, GREEN, FONT_SB)])
        s.shapes.add_picture(str(pose), Inches(px), Inches(y + 0.28),
                             width=Inches(img_w))
        ih = img_w * 450 / 800
        add_text(s, px, y + 0.28 + ih + 0.10, img_w, 0.90,
                 [P("Body pose recovered for students with no detectable face. The single "
                    "upright skeleton is the standing teacher — correctly distinguished "
                    "from every seated student, which is a real check that the signal "
                    "means something.", 8.5, False, MUTE, line=1.18)])

    yy = y + 3.54
    small = [("PROBLEM 3", "Back-row students missed entirely",
              "Inference resolution 960 px shrank a 60 px student to ~30 px.",
              "imgsz 1280 + conf 0.30 → 139 to 236 persons."),
             ("PROBLEM 4", "Low-resolution input collapses face detection",
              "Measured 7% face rate below 480 px vs 44% at 720–1080 px.",
              "Set a minimum capture resolution as a deployment requirement.")]
    cwid2 = (CW - 0.28) / 2
    for i, (tag, title, cause, fix) in enumerate(small):
        x = ML + i * (cwid2 + 0.28)
        rect(s, x, yy, cwid2, 1.42, fill=PANEL, line=BORDER)
        add_text(s, x + 0.22, yy + 0.16, cwid2 - 0.44, 0.22,
                 [P(tag, 8.5, True, AMBER, FONT_SB)])
        add_text(s, x + 0.22, yy + 0.40, cwid2 - 0.44, 0.30,
                 [P(title, 10.5, True, INK, FONT_SB, line=1.10)])
        add_text(s, x + 0.22, yy + 0.74, cwid2 - 0.44, 0.34,
                 [P(cause, 9, False, BODY, line=1.16)])
        add_text(s, x + 0.22, yy + 1.10, cwid2 - 0.44, 0.26,
                 [PR([R("→ ", 9, True, GREEN), R(fix, 9, False, TEAL_D)], line=1.16)])

    return s


def s16_problems_interpretation(prs, page):
    s, y = chrome(prs, "Criterion 5 — Problems Faced (2 / 2)",
                  "Interpretation problems — and where we say “we cannot tell”",
                  None, page)

    # Problem 5 - gaze inversion
    rect(s, ML, y, 6.05, 2.30, fill=WHITE, line=BORDER)
    bar(s, ML, y, 6.05, 0.05, RED)
    add_text(s, ML + 0.24, y + 0.22, 5.6, 0.24,
             [P("PROBLEM 5", 9, True, RED, FONT_SB)])
    add_text(s, ML + 0.24, y + 0.50, 5.6, 0.30,
             [P("Gaze direction was silently inverted", 13, True, INK, FONT_SB)])
    add_text(s, ML + 0.24, y + 0.90, 5.6, 0.80,
             [P("The head-pose model reports pitch as up-positive; our contract assumed "
                "down-positive. Students bowed over desks were being labelled “looking "
                "back/up”, and “down” could never fire at all. Confirmed by reading the "
                "library's own source, not guessed.", 10, False, BODY, line=1.20)])
    add_text(s, ML + 0.24, y + 1.76, 5.6, 0.44,
             [PR([R("Solution:  ", 10, True, GREEN, FONT_SB),
                  R("negate at the source + a regression test pinning the convention. "
                    "“down” labels went 0 → 6 on the same images; “back” 2 → 0.",
                    10, False, TEAL_D)], line=1.20)])

    # Problem 6 - the two distraction cases (the user's key ask)
    x2 = ML + 6.05 + 0.28
    w2 = SW - MR - x2
    rect(s, x2, y, w2, 2.30, fill=WHITE, line=BORDER)
    bar(s, x2, y, w2, 0.05, AMBER)
    add_text(s, x2 + 0.24, y + 0.22, w2 - 0.48, 0.24,
             [P("PROBLEM 6 — THE TWO DISTRACTION CASES", 9, True, AMBER, FONT_SB)])
    add_text(s, x2 + 0.24, y + 0.50, w2 - 0.48, 0.30,
             [P("“Head down” means two different things", 13, True, INK, FONT_SB)])
    sub = (w2 - 0.62) / 2
    cases = [("head_down_WITH_device", GREEN,
              "Gaze down + a phone detected overlapping the student. Defensible "
              "behavioural reading → we do flag this."),
             ("head_down_NO_device", AMBER,
              "Gaze down, nothing detected. Could be reading, writing, thinking — or "
              "disengaged. We refuse to guess.")]
    for i, (name, c, desc) in enumerate(cases):
        xx = x2 + 0.24 + i * (sub + 0.14)
        rect(s, xx, y + 0.88, sub, 1.30, fill=PANEL, line=BORDER)
        add_text(s, xx + 0.14, y + 1.00, sub - 0.28, 0.36,
                 [P(name, 8.5, True, c, MONO, line=1.10)])
        add_text(s, xx + 0.14, y + 1.40, sub - 0.28, 0.72,
                 [P(desc, 9, False, BODY, line=1.18)])

    # Bottom: the principle + peer interaction problem
    yy = y + 2.56
    rect(s, ML, yy, 6.05, 1.62, fill=INK, line=None)
    add_text(s, ML + 0.26, yy + 0.20, 5.55, 0.28,
             [P("THE PRINCIPLE WE ADOPTED", 9.5, True, RGBColor(0x6F, 0xC5, 0xCE), FONT_SB)])
    add_text(s, ML + 0.26, yy + 0.54, 5.55, 0.94,
             [P("Keep ambiguous cases ambiguous. BOSS — a validated human observation "
                "instrument — has the same unresolved confusion between “quiet, head down, "
                "working” and “quiet, head down, disengaged”. Forcing a confident split "
                "would be less honest than the established instrument itself.",
                10.5, False, WHITE, line=1.24)])

    card(s, x2, yy, w2, 1.62,
         heading="PROBLEM 7 — Peer interaction: a false positive we found and kept",
         lines=["Our geometric pair detector flagged its highest-confidence pair on real "
                "footage. We rendered it and looked: two students at different, "
                "non-adjacent desks, both bent over their own work, not interacting.",
                "We documented the false positive in the code and measured that the "
                "proximity threshold is looser than real desk spacing — rather than "
                "quietly tuning it until the demo looked good."],
         accent=AMBER, heading_size=11.5, body_size=9.5, fill=PANEL)
    return s


def s17_timeline(prs, page):
    s, y = chrome(prs, "Criterion 5 — Plan",
                  "Timeline and team split",
                  "Difficulty: moderate. Feasible on our own GPU using pretrained models — "
                  "no training from scratch is required for the core pipeline.", page)
    months = [
        ("M1", "Perception stack", "YOLOv11 · Face Mesh · head pose · body pose", True),
        ("M2", "Identity", "ByteTrack + CAL Re-ID, PRCC/LTCC benchmark", True),
        ("M3", "Scene graph", "graph construction + annotate own clips", False),
        ("M4", "Temporal + GAR", "STTran-style decoder, ARG readout", False),
        ("M5", "Analytics + report", "dashboard, evaluation, write-up", False),
    ]
    gap = 0.26
    bw = (CW - 4 * gap) / 5
    for i, (m, title, detail, done) in enumerate(months):
        x = ML + i * (bw + gap)
        rect(s, x, y, bw, 1.66, fill=WHITE if done else PANEL,
             line=GREEN if done else BORDER, line_w=1.4 if done else 0.75)
        add_text(s, x + 0.18, y + 0.18, 0.7, 0.26,
                 [P(m, 11.5, True, GREEN if done else MUTE, FONT_SB)])
        if done:
            pill(s, x + bw - 0.78, y + 0.18, 0.60, 0.24, "DONE", GREEN, WHITE, 7.5)
        add_text(s, x + 0.18, y + 0.54, bw - 0.36, 0.30,
                 [P(title, 11.5, True, INK, FONT_SB, line=1.08)])
        add_text(s, x + 0.18, y + 0.92, bw - 0.36, 0.64,
                 [P(detail, 8.5, False, BODY, line=1.20)])
        bar(s, x + 0.10, y + 1.80, bw - 0.20, 0.055,
            GREEN if done else RGBColor(0xD8, 0xE0, 0xE8))

    yy = y + 2.28
    add_text(s, ML, yy, CW, 0.3,
             [P("Team of three — vertical subsystem ownership", 13.5, True, INK, FONT_SB)])
    yy += 0.42
    members = [("A", "Perception & Identity", "Stages 1–2 · Re-ID benchmark · "
                "per-student trajectories"),
               ("B", "Relational & Temporal", "Stages 3–4 · scene graph · "
                "peer-influence modelling"),
               ("C", "Group Activity & App", "Stage 5 · dashboard · evaluation · "
                "predictive disengagement")]
    cwid = (CW - 2 * 0.26) / 3
    for i, (letter, role, detail) in enumerate(members):
        x = ML + i * (cwid + 0.26)
        rect(s, x, yy, cwid, 1.22, fill=PANEL, line=BORDER)
        rect(s, x + 0.20, yy + 0.22, 0.46, 0.46, fill=INK, line=None, adj=0.22)
        add_text(s, x + 0.20, yy + 0.30, 0.46, 0.3,
                 [P(letter, 13, True, WHITE, FONT_SB, align=PP_ALIGN.CENTER)])
        add_text(s, x + 0.80, yy + 0.22, cwid - 1.0, 0.28,
                 [P(role, 12, True, INK, FONT_SB)])
        add_text(s, x + 0.80, yy + 0.54, cwid - 1.0, 0.56,
                 [P(detail, 9.5, False, BODY, line=1.18)])

    card(s, ML, yy + 1.44, CW, 0.82,
         heading="Risk that is already retired",
         lines=["The single largest technical risk — “can this perception stack actually work "
                "on real classroom footage from a bad camera angle?” — has been answered "
                "empirically, with the failure modes found, measured and mitigated."],
         accent=GREEN, heading_size=11.5, body_size=10, fill=PANEL2)
    return s


def s18_novelty(prs, page):
    s, y = chrome(prs, "Originality",
                  "Three contributions, each requiring the whole system",
                  "None of these are possible with per-frame detection alone — that is the "
                  "argument for why all three pillars must coexist.", page)
    items = [
        ("01", "Identity-persistent engagement trajectories",
         "Because identity is stable for the session, we can plot attention per student "
         "across an entire lecture and locate the moment it drops — instead of emitting "
         "disconnected per-frame scores. The 15-second rolling window and per-student "
         "baseline are already implemented."),
        ("02", "Peer-orientation edges in the graph",
         "Student→student edges let us ask whether disengagement propagates through "
         "neighbours — an underexplored social signal. Grounded in Kendon's F-formations, "
         "and deliberately reported as “oriented toward each other”, never as "
         "“off-task”, because vision cannot judge the content of a conversation."),
        ("03", "Honest uncertainty as a first-class output",
         "The system reports which students it cannot assess and which behaviours it "
         "cannot disambiguate, rather than emitting a confident label anyway. The "
         "reviewed classroom literature consistently omits this, and it is what makes "
         "the output safe to put in front of a teacher."),
    ]
    cwid = (CW - 2 * 0.26) / 3
    for i, (num, title, body) in enumerate(items):
        x = ML + i * (cwid + 0.26)
        rect(s, x, y, cwid, 3.20, fill=WHITE, line=BORDER)
        add_text(s, x + 0.24, y + 0.24, cwid - 0.48, 0.5,
                 [P(num, 26, True, RGBColor(0xC9, 0xDA, 0xE2), FONT_SB)])
        add_text(s, x + 0.24, y + 0.84, cwid - 0.48, 0.68,
                 [P(title, 13, True, INK, FONT_SB, line=1.10)])
        add_text(s, x + 0.24, y + 1.60, cwid - 0.48, 1.44,
                 [P(body, 10, False, BODY, line=1.22)])

    yy = y + 3.46
    rect(s, ML, yy, CW, 1.02, fill=PANEL2, line=TEAL, line_w=1.25)
    add_text(s, ML + 0.28, yy + 0.18, CW - 0.56, 0.70,
             [PR([R("Framed honestly:  ", 11, True, TEAL_D, FONT_SB),
                  R("novel in combination and in domain — not a claim to have improved "
                    "Re-ID, scene-graph generation or group-activity recognition as "
                    "individual techniques. That is a defensible position for a course "
                    "project and an honest one for a paper.", 11, False, INK)], line=1.24)])
    return s


def s19_close(prs, page):
    s, y = chrome(prs, None, "Where we stand against each criterion", None, page)
    rows = [
        ("1", "Problem Understanding",
         "Behaviour, not attendance. Engagement defined as observable behaviour before "
         "it is measured, with six objectives each carrying a success criterion.", GREEN),
        ("2", "Research Paper Selection",
         "CAL (CVPR 2022) · STTran (ICCV 2021) · ARG (CVPR 2019) — one per pillar, "
         "top-tier, public code, correctly matched.", GREEN),
        ("3", "Literature Review Depth",
         "Four technical areas compared on method/benchmark/limitation, plus an "
         "education-psychology layer that determined our labels and our 15 s window.", GREEN),
        ("4", "Dataset Identification",
         "SCB-Dataset as classroom-native primary (13,330 images / 122,977 labels), "
         "field-standard benchmarks per module, and the fragmentation challenge stated "
         "with a mitigation.", GREEN),
        ("5", "Project Feasibility & Plan",
         "Stages 1–2 built and tested (109 tests, 100% schema-valid output), seven real "
         "problems found and fixed, 5-month plan with the biggest risk already retired.",
         GREEN),
    ]
    for i, (num, title, body, c) in enumerate(rows):
        yy = y + i * 0.92
        rect(s, ML, yy, CW, 0.80, fill=WHITE, line=BORDER)
        bar(s, ML, yy, 0.055, 0.80, c)
        rect(s, ML + 0.24, yy + 0.19, 0.42, 0.42, fill=INK, line=None, adj=0.22)
        add_text(s, ML + 0.24, yy + 0.25, 0.42, 0.3,
                 [P(num, 12.5, True, WHITE, FONT_SB, align=PP_ALIGN.CENTER)])
        add_text(s, ML + 0.82, yy + 0.14, 3.15, 0.3,
                 [P(title, 12.5, True, INK, FONT_SB)])
        add_text(s, ML + 0.82, yy + 0.42, CW - 1.10, 0.32,
                 [P(body, 9.5, False, BODY, line=1.14)])
        add_text(s, ML + 4.05, yy + 0.14, 0.3, 0.3, [P("", 9)])

    yy = y + 5 * 0.92 + 0.16
    rect(s, ML, yy, CW, 0.92, fill=INK, line=None)
    add_text(s, ML + 0.30, yy + 0.20, CW - 0.60, 0.56,
             [PR([R("What separates this review:  ", 11.5, True,
                    RGBColor(0x6F, 0xC5, 0xCE), FONT_SB),
                  R("the results slides report bugs we found in our own system and limits "
                    "we cannot engineer away — because a plan that survives contact with "
                    "real data is worth more than one that has never been tested.",
                    11.5, False, WHITE)], line=1.24)])
    return s


# --------------------------------------------------------------------------- #

def main():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    s01_title(prs)
    s02_roadmap(prs, 2)
    s03_problem(prs, 3)
    s04_motivation(prs, 4)
    s05_objectives(prs, 5)
    s06_papers(prs, 6)
    s07_lit_technical(prs, 7)
    s08_lit_behavioural(prs, 8)
    s09_gap(prs, 9)
    s10_dataset_primary(prs, 10)
    s11_dataset_support(prs, 11)
    s12_architecture(prs, 12)
    s13_status(prs, 13)
    s14_results(prs, 14)
    s15_problems_perception(prs, 15)
    s16_problems_interpretation(prs, 16)
    s17_timeline(prs, 17)
    s18_novelty(prs, 18)
    s19_close(prs, 19)

    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
