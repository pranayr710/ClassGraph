"""Audit + preview-render the generated deck.

No PowerPoint/LibreOffice on this machine, so layout bugs cannot be caught by
opening the file. This does two things instead:

1. Geometric audit -- shapes off-slide, and text likely to overflow its box
   (estimated from character count vs box area at the run's font size).
2. Approximate PIL render of each slide so the layout can actually be looked
   at. Text wrapping is approximated, so treat the render as a layout sanity
   check, not a pixel-accurate proof.

Run:  python verify_ppt.py
"""

from __future__ import annotations

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent
PPTX = ROOT / "ClassGraph.pptx"
PREVIEW = ROOT / "ppt_preview"
PREVIEW.mkdir(exist_ok=True)

SCALE = 110  # px per inch for the preview
EMU_IN = 914400


def emu_in(v) -> float:
    return (v or 0) / EMU_IN


def find_font(size_px: int, bold: bool = False):
    names = (["segoeuib.ttf", "seguisb.ttf", "arialbd.ttf"] if bold
             else ["segoeui.ttf", "arial.ttf"])
    for n in names:
        try:
            return ImageFont.truetype(n, size_px)
        except OSError:
            continue
    return ImageFont.load_default()


def iter_runs(shape):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip():
                yield para, run


def audit(prs) -> list[str]:
    sw, sh = emu_in(prs.slide_width), emu_in(prs.slide_height)
    issues: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            x, y = emu_in(shape.left), emu_in(shape.top)
            w, h = emu_in(shape.width), emu_in(shape.height)
            # 1. off-slide
            if x < -0.02 or y < -0.02:
                issues.append(f"slide {idx}: shape starts off-slide at ({x:.2f},{y:.2f})")
            if x + w > sw + 0.02:
                issues.append(
                    f"slide {idx}: shape overflows RIGHT edge by "
                    f"{x + w - sw:.2f}in (x={x:.2f} w={w:.2f})")
            if y + h > sh + 0.02:
                issues.append(
                    f"slide {idx}: shape overflows BOTTOM edge by "
                    f"{y + h - sh:.2f}in (y={y:.2f} h={h:.2f})")
            # 2. text overflow estimate
            if shape.has_text_frame and w > 0.2:
                # Accumulate per-run so a 27pt heading above 9pt body text does
                # not inflate every line to 27pt (which over-reported badly).
                needed = 0.0
                for para, run in iter_runs(shape):
                    size_pt = run.font.size.pt if run.font.size else 12.0
                    # avg glyph advance ~0.50em for Segoe UI body text
                    chars_per_line = max(1.0, (w * 72.0) / (0.50 * size_pt))
                    lines = math.ceil(len(run.text) / chars_per_line)
                    needed += lines * size_pt * 1.26 / 72.0
                if needed:
                    if needed > h + 0.10:
                        issues.append(
                            f"slide {idx}: text may overflow box "
                            f"(need ~{needed:.2f}in, have {h:.2f}in) :: "
                            f"{shape.text_frame.text[:58]!r}")
    return issues


def render(prs) -> None:
    sw, sh = emu_in(prs.slide_width), emu_in(prs.slide_height)
    W, H = int(sw * SCALE), int(sh * SCALE)
    for idx, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (W, H), (255, 255, 255))
        d = ImageDraw.Draw(img)
        for shape in slide.shapes:
            x, y = emu_in(shape.left) * SCALE, emu_in(shape.top) * SCALE
            w, h = emu_in(shape.width) * SCALE, emu_in(shape.height) * SCALE

            # Tables are GraphicFrames, not autoshapes -- render them cell by
            # cell, otherwise five table-driven slides preview as blank.
            if getattr(shape, "has_table", False):
                tbl = shape.table
                col_x = [x]
                for col in tbl.columns:
                    col_x.append(col_x[-1] + emu_in(col.width) * SCALE)
                row_y = [y]
                for row in tbl.rows:
                    row_y.append(row_y[-1] + emu_in(row.height) * SCALE)
                for ri, row in enumerate(tbl.rows):
                    for ci in range(len(tbl.columns)):
                        cell = tbl.cell(ri, ci)
                        cx0, cx1 = col_x[ci], col_x[ci + 1]
                        cy0, cy1 = row_y[ri], row_y[ri + 1]
                        try:
                            cf = tuple(cell.fill.fore_color.rgb)
                        except Exception:
                            cf = (255, 255, 255)
                        d.rectangle([cx0, cy0, cx1, cy1], fill=cf,
                                    outline=(214, 224, 232))
                        para = cell.text_frame.paragraphs[0]
                        if not para.runs:
                            continue
                        r0 = para.runs[0]
                        sz = r0.font.size.pt if r0.font.size else 9.0
                        try:
                            col_rgb = tuple(r0.font.color.rgb)
                        except Exception:
                            col_rgb = (60, 80, 102)
                        spx = max(6, int(sz * SCALE / 72.0))
                        fnt = find_font(spx, bool(r0.font.bold))
                        avail = max(1, int((cx1 - cx0 - 12) / (spx * 0.52)))
                        ty = cy0 + 4
                        line = ""
                        for word in cell.text.replace("\n", " ").split():
                            if len(line) + len(word) + 1 <= avail:
                                line = (line + " " + word).strip()
                            else:
                                d.text((cx0 + 6, ty), line, fill=col_rgb, font=fnt)
                                ty += spx * 1.22
                                line = word
                        if line:
                            d.text((cx0 + 6, ty), line, fill=col_rgb, font=fnt)
                continue
            # fill
            fill_rgb = None
            try:
                if shape.fill.type is not None and shape.fill.type == 1:
                    fill_rgb = tuple(shape.fill.fore_color.rgb)
            except Exception:
                fill_rgb = None
            line_rgb = None
            try:
                if shape.line.color and shape.line.color.type is not None:
                    line_rgb = tuple(shape.line.color.rgb)
            except Exception:
                line_rgb = None
            if fill_rgb or line_rgb:
                d.rectangle([x, y, x + w, y + h], fill=fill_rgb, outline=line_rgb)
            if shape.shape_type == 13:  # picture
                d.rectangle([x, y, x + w, y + h], fill=(200, 210, 220),
                            outline=(120, 140, 160))
                d.text((x + 6, y + 6), "[image]", fill=(60, 80, 100),
                       font=find_font(12))
            # text
            if shape.has_text_frame:
                ty = y + 2
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs)
                    if not txt.strip():
                        continue
                    run0 = para.runs[0]
                    size_pt = run0.font.size.pt if run0.font.size else 12.0
                    bold = bool(run0.font.bold)
                    try:
                        col = tuple(run0.font.color.rgb)
                    except Exception:
                        col = (40, 60, 80)
                    size_px = max(7, int(size_pt * SCALE / 72.0))
                    font = find_font(size_px, bold)
                    # crude wrap
                    avg = max(1, int(w / (size_px * 0.52)))
                    line = ""
                    for word in txt.split():
                        if len(line) + len(word) + 1 <= avg:
                            line = (line + " " + word).strip()
                        else:
                            d.text((x + 3, ty), line, fill=col, font=font)
                            ty += size_px * 1.24
                            line = word
                    if line:
                        d.text((x + 3, ty), line, fill=col, font=font)
                        ty += size_px * 1.24
        d.rectangle([0, 0, W - 1, H - 1], outline=(190, 200, 210))
        img.save(PREVIEW / f"slide_{idx:02d}.png")


def main():
    prs = Presentation(str(PPTX))
    issues = audit(prs)
    print(f"=== AUDIT: {len(issues)} potential issues ===")
    for i in issues:
        print("  " + i)
    render(prs)
    print(f"\nrendered {len(prs.slides._sldIdLst)} previews -> {PREVIEW}")


if __name__ == "__main__":
    main()
